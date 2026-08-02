#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成 PSPP 入门教学讲义 PPTX。
基于 docs/tutorial/ 下的 10 章 Markdown 教程，自动抽取章节标题与要点，
输出到 docs/public/downloads/pspp-teaching.pptx。

用法:
    python tools/make_teaching_pptx.py
"""
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

TUTORIAL_DIR = r"D:\PSPP-docs\docs\tutorial"
OUT = r"D:\PSPP-docs\docs\public\downloads\pspp-teaching.pptx"

# 品牌色（与站点主题协调）
PRIMARY = RGBColor(0x1F, 0x7A, 0x8C)   # 青蓝
DARK = RGBColor(0x22, 0x2A, 0x35)      # 深灰
GREY = RGBColor(0x55, 0x5F, 0x6B)      # 中灰
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCF, 0xE8, 0xEC)     # 浅青
FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def md_clean(text: str) -> str:
    """去掉 markdown 行内格式，保留可读文本。"""
    text = re.sub(r"!\[[^\]]*\]\([^)]*\)", "", text)          # 图片
    text = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", text)       # 链接 -> 文字
    text = re.sub(r"`([^`]+)`", r"\1", text)                   # 行内代码
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)             # 粗体
    text = re.sub(r"\*([^*]+)\*", r"\1", text)                 # 斜体
    text = text.replace("\\", "")
    return text.strip()


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def text(slide, x, y, w, h, content, size=18, color=DARK, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = content
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT
    return tb


def title_slide():
    s = add_slide()
    bg(s, PRIMARY)
    rect(s, 0, Inches(3.55), SW, Inches(0.06), LIGHT)
    text(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.5),
         "PSPP 统计软件", size=54, color=WHITE, bold=True)
    text(s, Inches(0.9), Inches(3.75), Inches(11.5), Inches(1.0),
         "入门教学讲义", size=40, color=WHITE, bold=True)
    text(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.6),
         "GNU PSPP 2.0.1 · 中文手册与教程配套", size=20, color=LIGHT)
    text(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
         "PSPP 文档中心", size=16, color=LIGHT)


def divider_slide(num, title):
    s = add_slide()
    bg(s, DARK)
    rect(s, 0, Inches(3.25), SW, Inches(0.05), PRIMARY)
    text(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.0),
         f"第 {num} 章", size=24, color=PRIMARY, bold=True)
    text(s, Inches(0.9), Inches(2.95), Inches(11.5), Inches(1.3),
         title, size=40, color=WHITE, bold=True)


def content_slide(title, bullets):
    s = add_slide()
    bg(s, WHITE)
    rect(s, 0, 0, SW, Inches(0.16), PRIMARY)
    text(s, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.9),
         title, size=30, color=PRIMARY, bold=True)
    tb = s.shapes.add_textbox(Inches(0.85), Inches(1.65), Inches(11.6), Inches(5.4))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for lvl, b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(11)
        p.level = lvl
        r = p.add_run()
        r.text = ("• " if lvl == 0 else "– ") + b
        r.font.size = Pt(20 if lvl == 0 else 17)
        r.font.color.rgb = DARK if lvl == 0 else GREY
        r.font.name = FONT
    return s


def parse_chapters():
    files = sorted(
        f for f in os.listdir(TUTORIAL_DIR)
        if re.match(r"^\d{2}-.*\.md$", f)
    )
    chapters = []
    for fn in files:
        path = os.path.join(TUTORIAL_DIR, fn)
        with open(path, encoding="utf-8") as fh:
            lines = fh.read().splitlines()
        title = None
        sections = []          # (h2_title, [(level, bullet_text), ...])
        cur = None
        for ln in lines:
            m1 = re.match(r"^#\s+(.*)$", ln)
            if m1:
                title = md_clean(m1.group(1))
                continue
            m2 = re.match(r"^##\s+(.*)$", ln)
            if m2:
                cur = {"title": md_clean(m2.group(1)), "items": []}
                sections.append(cur)
                continue
            m3 = re.match(r"^###\s+(.*)$", ln)
            if m3 and cur is not None:
                cur["items"].append((1, md_clean(m3.group(1))))
                continue
            # 列表项
            mli = re.match(r"^\s*[-*]\s+(.*)$", ln)
            if mli and cur is not None:
                txt = md_clean(mli.group(1))
                if txt:
                    cur["items"].append((0, txt))
        if title:
            chapters.append((title, sections))
    return chapters


def build():
    chapters = parse_chapters()
    title_slide()

    # 大纲页
    s = add_slide()
    bg(s, WHITE)
    rect(s, 0, 0, SW, Inches(0.16), PRIMARY)
    text(s, Inches(0.6), Inches(0.45), Inches(12), Inches(0.9),
         "课程大纲", size=32, color=PRIMARY, bold=True)
    tb = s.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.5), Inches(5.4))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for i, (t, _) in enumerate(chapters, 1):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(7)
        r = p.add_run()
        r.text = f"{i}.  {t}"
        r.font.size = Pt(19)
        r.font.color.rgb = DARK
        r.font.name = FONT

    # 逐章
    for i, (title, sections) in enumerate(chapters, 1):
        divider_slide(i, title)
        if not sections:
            continue
        for sec in sections:
            items = sec["items"]
            if not items:
                continue
            # 每页最多 8 条；超出则拆分
            chunks = [items[s:s + 8] for s in range(0, len(items), 8)]
            for ci, chunk in enumerate(chunks):
                sec_title = sec["title"] if ci == 0 else f"{sec['title']}（续）"
                content_slide(sec_title, chunk)

    # 结束页
    s = add_slide()
    bg(s, PRIMARY)
    text(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.2),
         "参考与延伸", size=40, color=WHITE, bold=True)
    text(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(2.4),
         "• 中文参考手册 / 英文参考手册（站内）\n"
         "• GNU PSPP 2.0.1 官方手册（FDL 许可）\n"
         "• 本讲义与《PSPP 入门教程》内容一致，可配合 DOCX / PDF 使用",
         size=20, color=LIGHT)

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"OK: saved {OUT}")
    print(f"slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")


if __name__ == "__main__":
    build()
